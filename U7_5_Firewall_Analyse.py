#!/usr/bin/python3
##############################################
#
# Name: U7_5_Firewall_Analyse.py
#
# Author: Peter Christen
#
# Version: 1.1
#
# Date: 26.09.2017
#       27.09.2017: V1.1 Div. Anpassungen
#
# Purpose: Analysiert Firewall Log Daten
#
# Hinweis: Braucht die Module matplotlib und python-pptx
#
##############################################

# Module
import sqlite3
import openpyxl
from openpyxl import load_workbook
from openpyxl.styles import PatternFill, Border, Side, Alignment, Protection, Font, Color, colors
import datetime

# Variablen
connection = sqlite3.connect(':memory:')
cursor = connection.cursor()
ExcelName = "U7_5_Firewall_Log_Analyse.xlsx"
Titel = 'Firewall Analyse'
now = datetime.datetime.now()
today = now.strftime("%d.%m.%Y")

sid = {
    '2501': 'DB1',
    '2502': 'DB2',
    '2511': 'DB3',
    '2512': 'DB4',
    '2513': 'DB5',
    '2514': 'DB6',
    '2525': 'DB7',
    '2526': 'DB8',
    '2527': 'DB9',
    '2528': 'DB10',
    '2529': 'DB11',
    '2530': 'DB12',
    '3024': 'DB13',
    '3016': 'DB14',
    '636': 'DB15'}

################### Keine Aenderungen mehr nötig ab hier ################

# Functions


def colsize(col):
    if col > 90:
        col = col - 26
        b = chr(col)
        col = 'A' + b
    else:
        col = str(chr(col))
    return col


def createdb():
    # Datenbank erstellen
    cursor.execute(
        'create table if not exists firelog ( sourceip varchar(20), targetip varchar(20), port varchar(20), sid varchar(10), count int)')
    cursor.execute(
        'CREATE INDEX if not exists firelog_ind on firelog (sourceip, targetip, port)')


def read_excel():
    # Firewall log einlesen

    # xlsx-File öffnen
    wb = load_workbook(
        filename='U7_5_Firewall_Log_Auszug.xlsx',
        read_only=True)
    sheet1 = wb.worksheets[0]
    worksheet = wb[sheet1.title]

    # xlsx-File einlesen
    r = 0
    w = []
    for row in worksheet.iter_rows():
        r += 1
        c = 0
        for cell in row:
            c += 1
            w.append(cell.value)

        if r > 2:
            po = str(w[2]).split(".")
            cursor.execute("replace into firelog values(?,?,?,?,?)",
                           (w[0], w[1], po[0], sid[po[0]], w[4]))
        del w[:]

    cursor.execute("delete from firelog where sourceip like 'Source IP'")
    connection.commit()


def write_excel():
    # Erstelle Excel aus Firewall DB

    wb = openpyxl.Workbook()
    ws1 = wb.worksheets[0]
    ws1.title = 'Analyse'

    # Styles definieren
    # Font
    fontT = Font(bold=True, size=14, color=colors.BLACK)
    fontb = Font(bold=True, color=colors.BLACK)

    # Alignment
    alignC = Alignment(
        horizontal='center',
        vertical='top',
        text_rotation=0,
        shrink_to_fit=False,
        wrap_text=False)
    alignR = Alignment(
        horizontal='right',
        vertical='top',
        text_rotation=0,
        shrink_to_fit=False,
        wrap_text=False)
    alignL = Alignment(
        horizontal='left',
        vertical='top',
        text_rotation=0,
        shrink_to_fit=False,
        wrap_text=False)

    # Zellenfarben
    fillT = PatternFill(
        fill_type='solid',
        start_color='6EB7FF',
        end_color='6EB7FF')
    fillR = PatternFill(
        fill_type='solid',
        start_color='FF0000',
        end_color='FF0000')
    fillO = PatternFill(
        fill_type='solid',
        start_color='FCD020',
        end_color='FCD020')
    fillG = PatternFill(
        fill_type='solid',
        start_color='22B604',
        end_color='22B604')
    fillY = PatternFill(
        fill_type='solid',
        start_color='FFF251',
        end_color='FFF251')
    fillGh = PatternFill(
        fill_type='solid',
        start_color='e6ffcc',
        end_color='e6ffcc')
    fillOh = PatternFill(
        fill_type='solid',
        start_color='ffe6cc',
        end_color='ffe6cc')

    # Kollonengroesse definieren
    ws1.column_dimensions["A"].width = 14.0
    ws1.column_dimensions["B"].width = 14.0
    ws1.column_dimensions["C"].width = 8.0
    ws1.column_dimensions["D"].width = 7.0
    ws1.column_dimensions["E"].width = 7.0

    ws1.merge_cells('A1:E1')
    ws1['A1'].font = fontT
    ws1['A1'].fill = fillT
    ws1['A1'].value = Titel
    ws1['A3'].font = fontb
    ws1['A3'].value = "SourceIP"
    ws1['B3'].font = fontb
    ws1['B3'].value = "TargetIP"
    ws1['C3'].font = fontb
    ws1['C3'].value = "Port"
    ws1['D3'].font = fontb
    ws1['D3'].value = "SID"
    ws1['E3'].font = fontb
    ws1['E3'].value = "Hits"

    # Daten aus der Datenbanken einfuegen
    coln = 65
    z = 4

    cursor.execute("select * from firelog")

    for row in cursor:
        col = colsize(coln)
        ce = col + str(z)
        ws1[ce].value = row[0]
        col = colsize(coln + 1)
        ce = col + str(z)
        ws1[ce].value = row[1]
        col = colsize(coln + 2)
        ce = col + str(z)
        ws1[ce].value = row[2]
        col = colsize(coln + 3)
        ce = col + str(z)
        ws1[ce].value = row[3]
        col = colsize(coln + 4)
        ce = col + str(z)
        if row[4] > 500000:
            ws1[ce].fill = fillR
        elif row[4] > 100000:
            ws1[ce].fill = fillO
        else:
            ws1[ce].fill = fillG
        ws1[ce].value = row[4]
        z += 1

    # Summe der Hit Counts ausgeben
    dbnamen = []
    dbhits = []
    z = 3
    col = colsize(coln + 6)
    ce = col + str(z)
    ws1.merge_cells('G3:H3')
    ws1[ce].fill = fillT
    ws1[ce].font = fontb
    ws1[ce].value = "Summe der Hits"
    z += 1

    rows = cursor.execute(
        "select sid,sum(count) from firelog group by sid order by sum(count) desc")
    for row in rows:
        dbnamen.append(row[0])
        dbhits.append(row[1])
        col = colsize(coln + 6)
        ce = col + str(z)
        ws1[ce].value = row[0]
        col = colsize(coln + 7)
        ce = col + str(z)
        ws1[ce].value = row[1]
        z += 1

    wb.save(filename=ExcelName)
    create_grafik_pie(dbnamen, dbhits)
    create_presi(dbnamen, dbhits)


def create_grafik_pie(dbnamen, dbhits):
    # Module matplotlib
    import matplotlib.pyplot as plt

    fig1, ax1 = plt.subplots()
    ax1.pie(dbhits[0:10], labels=dbnamen[0:10],
            autopct='%1.1f%%', shadow=True, startangle=0)
    ax1.axis('equal')
    plt.title('Verteilung der Hits der 10 top Datenbank')

    fig1.savefig('DBHits.png')

    # plt.show()


def create_presi(dbnamen, dbhits):
    # Modul python-pptx
    from pptx import Presentation
    from pptx.util import Inches

    N = len(dbnamen)

    # Präsentation eröffnen
    prs = Presentation()

    # Slides definieren
    title_slide_layout = prs.slide_layouts[0]
    setting_slide_layout = prs.slide_layouts[1]

    # Slides hinzufügen
    slide0 = prs.slides.add_slide(title_slide_layout)
    slide1 = prs.slides.add_slide(setting_slide_layout)

    # Header Slide
    title = slide0.shapes.title
    subtitle = slide0.placeholders[1]
    title.text = "Verteilung Netzwerk Traffic"
    subtitle.text = "auf die 10 Top Datenbanken\n" + today

    # Setting Slide
    shapes = slide1.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Verteilung'

    tf = body_shape.text_frame
    beschreibung = 'Zehn von insgesamt ' + str(N) + ' Datenbanken'
    tf.text = beschreibung

    left = Inches(1.5)
    top = Inches(2.5)
    height = Inches(5)
    pic = slide1.shapes.add_picture('DBHits.png', left, top, height=height)

    # Präsentation sichern
    prs.save('U7_5_Firewall_Analyse_Presi.pptx')

# End Functions


# Main
createdb()
read_excel()
write_excel()

# End
cursor.close()
